import fitz  # PyMuPDF
import base64
from string import Template
from docx import Document

from langchain_core.messages import HumanMessage
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.output_parsers import BaseOutputParser

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

MODEL = "gemini-2.5-flash"

class QuizParser(BaseOutputParser):
    def parse(self, text: str):
        return text.content if hasattr(text, 'content') else text

def get_pdf_content(pdf_file):
    """Convert PDF pages to base64 images for multimodal understanding"""
    content_parts = []
    
    # Save uploaded file temporarily to read with fitz
    # Streamlit UploadedFile objects need to be handled carefully
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x extract so it's readable
        img_data = pix.tobytes("jpeg")
        
        base64_image = base64.b64encode(img_data).decode('utf-8')
        content_parts.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
        })
        
    return content_parts

def get_docx_text(docx_file):
    """Extract text from a Word (.docx) file"""
    try:
        
        text = ""
        doc = Document(docx_file)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except ImportError:
        return "Error: python-docx library not found. Please install it using 'pip install python-docx'"

def get_txt_text(txt_file):
    """Extract text from a TXT file"""
    try:
        return txt_file.read().decode('utf-8')
    except UnicodeDecodeError:
        txt_file.seek(0)
        return txt_file.read().decode('latin-1')

def get_pptx_content(pptx_file):
    """Extract text and images from a PowerPoint (.pptx) file"""
    try:
        content_parts = []
        prs = Presentation(pptx_file)
        
        for i, slide in enumerate(prs.slides):
            slide_text = f"\n--- Slide {i+1} ---\n"
            slide_images = []
            
            for shape in slide.shapes:
                # Extract Text
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                
                # Extract Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        # Support common image types
                        if image.content_type in ['image/jpeg', 'image/png', 'image/jpg']:
                            img_bytes = image.blob
                            base64_image = base64.b64encode(img_bytes).decode('utf-8')
                            slide_images.append({
                                "type": "image_url",
                                "image_url": {"url": f"data:{image.content_type};base64,{base64_image}"}
                            })
                    except Exception:
                        pass # Skip images that can't be processed

            # Add text part for the slide
            content_parts.append({"type": "text", "text": slide_text})
            # Add image parts for the slide
            content_parts.extend(slide_images)
            
        return content_parts
    except ImportError:
        return [{"type": "text", "text": "Error: python-pptx library not found. Please install it using 'pip install python-pptx'"}]
    except Exception as e:
        return [{"type": "text", "text": f"Error reading PPTX: {str(e)}"}]

def get_document_content(uploaded_files):
    """Extract content from multiple uploaded files (PDF->Images, DOCX/TXT/PPTX->Text+Images)"""
    content = []
    
    for file in uploaded_files:
        file.seek(0) # Ensure we read from start
        file_extension = file.name.lower().split('.')[-1]
        
        if file_extension == 'pdf':
            content.extend(get_pdf_content(file))
        elif file_extension == 'docx':
            text = get_docx_text(file)
            content.append({"type": "text", "text": f"\n--- Start of {file.name} ---\n{text}\n--- End of {file.name} ---\n"})
        elif file_extension == 'txt':
            text = get_txt_text(file)
            content.append({"type": "text", "text": f"\n--- Start of {file.name} ---\n{text}\n--- End of {file.name} ---\n"})
        elif file_extension == 'pptx':
            content.extend(get_pptx_content(file))
        else:
            print(f"Unsupported file type: {file.name}")
            
    return content

def get_quiz_chain(api_key, quiz_type="Classic", language="English"):
    try:
        if not api_key:
            return None
            
        llm = ChatGoogleGenerativeAI(google_api_key=api_key, model=MODEL, temperature=0.7)
        
        # Base instructions
        if quiz_type == "Test":
            instructions = Template("""
            You are an expert quiz maker. 
            Analyze the provided document images/text and create a multiple-choice quiz with $number questions.
            The questions should be appropriate for $level level students.
            
            Return the output STRICTLY as a VALID JSON ARRAY of objects.
            Each object must have the following structure:
            {
                "question": "The question text",
                "options": ["Option A", "Option B", "Option C", "Option D", "Option E"],
                "correct_answer": "The correct option text (must match one of the options exactly)"
            }

            The quiz and all questions/answers must be generated in the $language language.
            
            IMPORTANT:
            1. Return ONLY the JSON array.
            2. Do not include markdown formatting like ```json or ```.
            3. Ensure the JSON is valid and can be parsed.
            """)
        else:
            instructions = Template("""
            You are an expert quiz maker. 
            Analyze the provided document images/text and create a quiz with $number questions.
            The questions should be appropriate for $level level students.
            The quiz and all questions/answers must be generated in the $language language.
            
            Format the output as a numbered list of questions, followed by a numbered list of answers at the very end.
            """)

        # We return a function that constructs the full multimodal message
        def generate_quiz(inputs):
            content_parts = inputs['content'] 
            # Prepend instructions as text
            filled_instructions = instructions.substitute(
                number=inputs['number'], 
                level=inputs['level'],
                language=language
            )

            final_message_content = [{"type": "text", "text": filled_instructions}] + content_parts
            
            msg = HumanMessage(content=final_message_content)
            response = llm.invoke([msg])
            return response.content

        return generate_quiz
        
    except Exception as e:
        print(f"Error creating chain: {e}")
        return None

def get_abstract_chain(api_key, language="English"):
    try:
        if not api_key:
            return None
            
        llm = ChatGoogleGenerativeAI(google_api_key=api_key, model=MODEL, temperature=0.5)
        
        instructions = f"""
        You are an expert at summarizing documents.
        Create a concise abstract of the following content in {language}.
        The abstract should capture the main points and key ideas.
        """

        def generate_abstract(inputs):
            content_parts = inputs['content']
            final_message_content = [{"type": "text", "text": instructions}] + content_parts
            msg = HumanMessage(content=final_message_content)
            response = llm.invoke([msg])
            return response.content

        return generate_abstract
    except Exception as e:
        print(f"Error creating abstract chain: {e}")
        return None

def format_quiz_for_download(quiz_data, quiz_type):
    """Formats quiz data into a user-friendly printable string."""
    import json
    
    if quiz_type == "Test":
        try:
            # Try to parse if it's a string, otherwise assume it's already a list/dict
            if isinstance(quiz_data, str):
                # Clean up potential markdown
                cleaned_data = quiz_data.replace("```json", "").replace("```", "").strip()
                questions = json.loads(cleaned_data)
            else:
                questions = quiz_data
                
            output = []
            output.append("QUIZ GENERATED BY AI")
            output.append("====================")
            output.append("")
            
            # Questions Section
            for i, q in enumerate(questions):
                output.append(f"{i+1}. {q['question']}")
                for idx, opt in enumerate(q['options']):
                    letter = chr(97 + idx)  # 97 is 'a'
                    output.append(f"   {letter}) {opt}")
                output.append("")
            
            output.append("")
            output.append("ANSWER KEY")
            output.append("==========")
            
            # Answers Section
            for i, q in enumerate(questions):
                correct_text = q['correct_answer']
                try:
                    # Find index of correct answer in options to get the letter
                    correct_idx = q['options'].index(correct_text)
                    correct_letter = chr(97 + correct_idx)
                    output.append(f"{i+1}. {correct_letter}")
                except ValueError:
                    # Fallback if content doesn't match for some reason
                    output.append(f"{i+1}. {correct_text}")
                
            return "\n".join(output)
            
        except Exception:
            # Fallback if JSON parsing fails
            return str(quiz_data)
    else:
        # Classic mode is already text, just add a header
        return f"QUIZ GENERATED BY AI\n====================\n\n{str(quiz_data)}"
